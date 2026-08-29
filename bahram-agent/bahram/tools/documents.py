"""Document extraction for Bahram Agent."""

from __future__ import annotations

import logging
from pathlib import Path
from typing import Any, Optional

logger = logging.getLogger(__name__)


class DocumentExtractor:
    """Extract text from various document formats."""

    SUPPORTED_FORMATS = [
        ".pdf", ".docx", ".xlsx", ".pptx",
        ".txt", ".md", ".csv", ".json", ".yaml", ".yml",
        ".html", ".xml", ".rtf",
    ]

    def extract(self, filepath: str) -> dict[str, Any]:
        """Extract text from a document.

        Returns:
            Dict with 'content' and 'metadata' keys.
        """
        path = Path(filepath)
        if not path.exists():
            return {"error": f"File not found: {filepath}"}

        suffix = path.suffix.lower()

        try:
            if suffix == ".pdf":
                return self._extract_pdf(path)
            elif suffix == ".docx":
                return self._extract_docx(path)
            elif suffix == ".xlsx":
                return self._extract_xlsx(path)
            elif suffix == ".pptx":
                return self._extract_pptx(path)
            elif suffix in [".txt", ".md", ".csv", ".json", ".yaml", ".yml"]:
                return self._extract_text(path)
            elif suffix == ".html":
                return self._extract_html(path)
            elif suffix == ".xml":
                return self._extract_xml(path)
            else:
                return {"error": f"Unsupported format: {suffix}"}
        except ImportError as e:
            return {"error": f"Missing dependency: {e}"}
        except Exception as e:
            return {"error": f"Extraction failed: {e}"}

    def _extract_pdf(self, path: Path) -> dict:
        """Extract from PDF."""
        try:
            import PyPDF2
            reader = PyPDF2.PdfReader(str(path))
            content = "\n\n".join(
                page.extract_text() or "" for page in reader.pages
            )
            return {
                "content": content,
                "metadata": {
                    "pages": len(reader.pages),
                    "format": "pdf",
                },
            }
        except ImportError:
            return {"error": "PyPDF2 not installed. Run: pip install PyPDF2"}

    def _extract_docx(self, path: Path) -> dict:
        """Extract from DOCX."""
        try:
            from docx import Document
            doc = Document(str(path))
            content = "\n\n".join(
                para.text for para in doc.paragraphs if para.text
            )
            return {
                "content": content,
                "metadata": {
                    "paragraphs": len(doc.paragraphs),
                    "format": "docx",
                },
            }
        except ImportError:
            return {"error": "python-docx not installed. Run: pip install python-docx"}

    def _extract_xlsx(self, path: Path) -> dict:
        """Extract from XLSX."""
        try:
            from openpyxl import load_workbook
            wb = load_workbook(str(path), read_only=True)
            content_parts = []
            for sheet in wb.sheetnames:
                ws = wb[sheet]
                rows = []
                for row in ws.iter_rows(values_only=True):
                    rows.append("\t".join(str(cell or "") for cell in row))
                content_parts.append(f"## Sheet: {sheet}\n\n" + "\n".join(rows))
            wb.close()
            return {
                "content": "\n\n".join(content_parts),
                "metadata": {
                    "sheets": len(wb.sheetnames),
                    "format": "xlsx",
                },
            }
        except ImportError:
            return {"error": "openpyxl not installed. Run: pip install openpyxl"}

    def _extract_pptx(self, path: Path) -> dict:
        """Extract from PPTX."""
        try:
            from pptx import Presentation
            prs = Presentation(str(path))
            content_parts = []
            for i, slide in enumerate(prs.slides, 1):
                texts = []
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        texts.append(shape.text_frame.text)
                if texts:
                    content_parts.append(f"## Slide {i}\n\n" + "\n\n".join(texts))
            return {
                "content": "\n\n".join(content_parts),
                "metadata": {
                    "slides": len(prs.slides),
                    "format": "pptx",
                },
            }
        except ImportError:
            return {"error": "python-pptx not installed. Run: pip install python-pptx"}

    def _extract_text(self, path: Path) -> dict:
        """Extract from plain text files."""
        content = path.read_text(encoding="utf-8", errors="replace")
        return {
            "content": content,
            "metadata": {
                "format": path.suffix[1:],
                "size": len(content),
            },
        }

    def _extract_html(self, path: Path) -> dict:
        """Extract from HTML."""
        try:
            from bs4 import BeautifulSoup
            raw = path.read_text(encoding="utf-8", errors="replace")
            soup = BeautifulSoup(raw, "html.parser")
            content = soup.get_text(separator="\n", strip=True)
            return {
                "content": content,
                "metadata": {
                    "format": "html",
                    "title": soup.title.string if soup.title else None,
                },
            }
        except ImportError:
            # Fallback: just read raw
            content = path.read_text(encoding="utf-8", errors="replace")
            return {
                "content": content,
                "metadata": {"format": "html", "raw": True},
            }

    def _extract_xml(self, path: Path) -> dict:
        """Extract from XML."""
        content = path.read_text(encoding="utf-8", errors="replace")
        return {
            "content": content,
            "metadata": {"format": "xml"},
        }

    def get_supported_formats(self) -> list[str]:
        """Get list of supported formats."""
        return self.SUPPORTED_FORMATS
